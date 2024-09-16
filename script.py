import os
import argparse

import pydicom
import cv2
from cv2 import VideoWriter_fourcc
from pptx import Presentation
from pptx.util import Inches


def extract_images_from_dicom(dicom_seq_folder):
    images = [
        pydicom.dcmread(os.path.join(dicom_seq_folder, dicom_file))
        for dicom_file in os.listdir(dicom_seq_folder)
    ]
    images = [im.pixel_array for im in images if hasattr(im, "pixel_array")]
    maxes = max([im.max() for im in images]) if images else 1
    images = [
        im if maxes < 2**8 else im // (2**4) if maxes < 2**12 else im // 2**8 + 2**7
        for im in images
    ]
    images = [im.astype("uint8") for im in images]
    return images


def add_video_slide(prs, video_file):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(1)
    slide.shapes.add_movie(
        video_file,
        left,
        top,
        width=Inches(8.5),
        height=Inches(6),
        poster_frame_image=None,
        mime_type="video/mp4",
    )


def add_image_slide(prs, image_file):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(1)
    slide.shapes.add_picture(image_file, left, top, width=Inches(8.5), height=Inches(6))


def export_to_ppt(input_files, ppt_file):
    prs = Presentation()
    for input_file in input_files:
        if ".mp4" in input_file:
            add_video_slide(prs, input_file)
        else:
            add_image_slide(prs, input_file)
    prs.save(ppt_file)


def setup_parser():
    parser = argparse.ArgumentParser(
        prog="DICOM2PPT",
        description="Converts multiple DICOM files to images, GIFs or mp4 videos and insert them into a PPT file",
        epilog="Usage: inputfolder outputfolder",
    )
    parser.add_argument("inputfolder")
    parser.add_argument("outputfolder")
    return parser


def main():
    parser = setup_parser()
    args = parser.parse_args()

    dicom_folder = args.inputfolder
    outpath = args.outputfolder

    media_paths = []

    for sequence_folder in os.listdir(dicom_folder):
        print(f"Reading {sequence_folder}")
        infolder = os.path.join(dicom_folder, sequence_folder)
        sequence = extract_images_from_dicom(infolder)
        if not sequence:
            print("Skipping...")
            continue
        first_frame_shape = sequence[0].shape[1::-1]
        format, codec = ".mp4", "mp4v"
        fourcc = VideoWriter_fourcc(*codec)
        fps = 1.0
        output_media_path = os.path.join(outpath, sequence_folder) + format
        output_writer = cv2.VideoWriter(
            output_media_path, fourcc, fps, first_frame_shape
        )
        assert output_writer.isOpened()
        for i, frame in enumerate(sequence):
            try:
                # imageio.mimsave(output_media_path, sequence)
                output_writer.write(frame)
            except Exception as e:
                # imageio.mimsave(output_media_path, [im for im in sequence if im.shape == sequence[0].shape])
                # imageio.imwrite(f"{output_media_path}.{i}.png", im)
                print(f"Exception frame ({i}): {e} \n\t{frame.max(), frame.min()}")

        media_paths.append(output_media_path)
        output_writer.release()
    export_to_ppt(media_paths, os.path.join(outpath, "presentation.pptx"))
    # cv2.destroyAllWindows()


if __name__ == "__main__":
    main()
